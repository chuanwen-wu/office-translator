#!/usr/bin/env python

import argparse
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
import json
import requests
from pptx.util import Pt

LANGUAGE_CODE_TO_LANGUAGE_ID = {
"""
Dict that maps Amazon Translate language code to MSO_LANGUAGE_ID enum value.

- Amazon Translate language codes: https://docs.aws.amazon.com/translate/latest/dg/what-is.html#what-is-languages
- python-pptx MSO_LANGUAGE_ID enum: https://python-pptx.readthedocs.io/en/latest/api/enum/MsoLanguageId.html

python-pptx doesn't support:
    - Azerbaijani (az)
    - Persian (fa)
    - Dari (fa-AF)
    - Tagalog (tl)
"""
    'af': MSO_LANGUAGE_ID.AFRIKAANS,
    'am': MSO_LANGUAGE_ID.AMHARIC,
    'ar': MSO_LANGUAGE_ID.ARABIC,
    'bg': MSO_LANGUAGE_ID.BULGARIAN,
    'bn': MSO_LANGUAGE_ID.BENGALI,
    'bs': MSO_LANGUAGE_ID.BOSNIAN,
    'cs': MSO_LANGUAGE_ID.CZECH,
    'da': MSO_LANGUAGE_ID.DANISH,
    'de': MSO_LANGUAGE_ID.GERMAN,
    'el': MSO_LANGUAGE_ID.GREEK,
    'en': MSO_LANGUAGE_ID.ENGLISH_US,
    'es': MSO_LANGUAGE_ID.SPANISH,
    'et': MSO_LANGUAGE_ID.ESTONIAN,
    'fi': MSO_LANGUAGE_ID.FINNISH,
    'fr': MSO_LANGUAGE_ID.FRENCH,
    'fr-CA': MSO_LANGUAGE_ID.FRENCH_CANADIAN,
    'ha': MSO_LANGUAGE_ID.HAUSA,
    'he': MSO_LANGUAGE_ID.HEBREW,
    'hi': MSO_LANGUAGE_ID.HINDI,
    'hr': MSO_LANGUAGE_ID.CROATIAN,
    'hu': MSO_LANGUAGE_ID.HUNGARIAN,
    'id': MSO_LANGUAGE_ID.INDONESIAN,
    'it': MSO_LANGUAGE_ID.ITALIAN,
    'ja': MSO_LANGUAGE_ID.JAPANESE,
    'ka': MSO_LANGUAGE_ID.GEORGIAN,
    'ko': MSO_LANGUAGE_ID.KOREAN,
    'lv': MSO_LANGUAGE_ID.LATVIAN,
    'ms': MSO_LANGUAGE_ID.MALAYSIAN,
    'nl': MSO_LANGUAGE_ID.DUTCH,
    'no': MSO_LANGUAGE_ID.NORWEGIAN_BOKMOL,
    'pl': MSO_LANGUAGE_ID.POLISH,
    'ps': MSO_LANGUAGE_ID.PASHTO,
    'pt': MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE,
    'ro': MSO_LANGUAGE_ID.ROMANIAN,
    'ru': MSO_LANGUAGE_ID.RUSSIAN,
    'sk': MSO_LANGUAGE_ID.SLOVAK,
    'sl': MSO_LANGUAGE_ID.SLOVENIAN,
    'so': MSO_LANGUAGE_ID.SOMALI,
    'sq': MSO_LANGUAGE_ID.ALBANIAN,
    'sr': MSO_LANGUAGE_ID.SERBIAN_LATIN,
    'sv': MSO_LANGUAGE_ID.SWEDISH,
    'sw': MSO_LANGUAGE_ID.SWAHILI,
    'ta': MSO_LANGUAGE_ID.TAMIL,
    'th': MSO_LANGUAGE_ID.THAI,
    'tr': MSO_LANGUAGE_ID.TURKISH,
    'uk': MSO_LANGUAGE_ID.UKRAINIAN,
    'ur': MSO_LANGUAGE_ID.URDU,
    'vi': MSO_LANGUAGE_ID.VIETNAMESE,
    'zh': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    'zh-TW': MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
}

# TERMINOLOGY_NAME = 'pptx-translator-terminology'
DONT_TRANSLATE_WORDS_FILE = f"./dont-translate-word.txt"
DONT_TRANSLATE_WORDS = f""

def contain_chinese(string):
    """
    检查整个字符串是否包含中文
    :param string: 需要检查的字符串
    :return: bool
    """
    for ch in string:
        if u'\u4e00' <= ch <= u'\u9fff':
            return True
 
    return False


def get_dont_translate_words(file_path):
    global DONT_TRANSLATE_WORDS
    arr = []
    with open(file_path) as f:
        for l in f.readlines():
            if l.strip().startswith('#'):
                continue
            arr.append(l.strip())
    
    DONT_TRANSLATE_WORDS = ",".join(arr)
    return ",".join(arr)


def translate_from_ollama(src, source_language_code, target_language_code):
    source_language = ""
    target_language = ""
    if source_language_code.lower() == 'en':
        source_language = "English"
    elif source_language_code.lower() == 'zh':
        source_language = "Chinese"

    if target_language_code.lower() == 'en':
        target_language = "English"
    elif target_language_code.lower() == 'zh':
        target_language = "Chinese"

    temp = {
        # "model": "llama3.2",
        "model": "qwen2.5",
        "stream": False,
        "options": {
            'temperature': 0
        },
        "messages": [
            {
                "role": "system",
                "content": f"You are a highly skilled translator. Your task is to accurately translate the text I provide from \
{source_language} into the {target_language} while preserving the meaning, tone, and nuance of the original text. \
Please maintain proper grammar, spelling, and punctuation in the translated version. \
Please answer briefly without any other hints or extensions. \
If the input text does not make sense or empty or incomplete, just return the original text. \
If the original text is already in English, return the original text directly. \
Remove the leading and trailing double quotes. \
Please keeping these phrases unchanged: {DONT_TRANSLATE_WORDS}."
            },
            {"role": "user", "content": f"{src}"}
        ]
    }
    # print(temp)

    data = json.dumps(temp)
    print(f'temp: {data}')
    # url = 'http://localhost:11434/api/generate'
    url = 'http://localhost:11434/api/chat'
    response = requests.post(url, data=data)
    # print(response.text)
    body = json.loads(response.text)
    # print(f"{src} -> {body['response']}")
    print(f"{src} --> {body['message']['content']}")
    return body['message']['content']


def delete_run(run):
    r = run._r
    r.getparent().remove(r)


def translate_text_frame(text_frame, source_language_code, target_language_code, terminology_names):
    for paragraph in text_frame.paragraphs:
        paraText = ''
        for index, paragraph_run in enumerate(paragraph.runs):
            paraText += paragraph_run.text
        # print(paraText)

        try:
            if len(paraText.strip()) == 0:
                continue

            # 源语言是中文，但实际文本未包含中文，跳过
            if contain_chinese(paraText.strip()) == False and source_language_code.lower() == 'zh':
                print(f"not chinese: {paraText}")
                continue

            reponse_text = translate_from_ollama(paraText.strip(), source_language_code, target_language_code)
            paragraph.runs[0].text = reponse_text
            size = len(paragraph.runs)
            index = size - 1
            while index > 0:
                # print(f"delete run {index} of {size}")
                delete_run(paragraph.runs[index])
                index = index - 1

            # # set font for translated text
            # if paragraph.font.size is not None:
            #     paragraph.font.size = paragraph.font.size - Pt(2)
            # elif paragraph.runs[0].font.size is not None:
            #     paragraph.runs[0].font.size = paragraph.runs[0].font.size - Pt(2)

                # for each in paragraph.runs: each.font.size = run.font.size
            
            # run.font.language_id = LANGUAGE_CODE_TO_LANGUAGE_ID[target_language_code]
            # paragraph.runs[index].text = reponse_text
            # paragraph.runs[index].font.language_id = LANGUAGE_CODE_TO_LANGUAGE_ID[target_language_code]
        except Exception as err:
            print(f"Unexpected {err=}, {type(err)=}")

def translate_table(table, source_language_code, target_language_code, terminology_names):
    for row in table.rows:
        for cell in row.cells:
            text_frame = cell.text_frame
            translate_text_frame(text_frame, source_language_code, target_language_code, terminology_names)
            

def translate_shape(shape, source_language_code, target_language_code, terminology_names):
    # table
    if shape.has_table:
        translate_table(shape.table, source_language_code, target_language_code, terminology_names)

    # text_frame
    if shape.has_text_frame:
        translate_text_frame(shape.text_frame, source_language_code, target_language_code, terminology_names)

    # groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            translate_shape(sub_shape, source_language_code, target_language_code, terminology_names)
    

def translate_presentation(presentation, source_language_code, target_language_code, terminology_names):
    slide_number = 1
    for slide in presentation.slides:
        print('Slide {slide_number} of {number_of_slides}'.format(
                slide_number=slide_number,
                number_of_slides=len(presentation.slides)))
        slide_number += 1

        # 尚不支持smartArt

        # translate comments
        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if len(text_frame.text) > 0:
                try:
                    print(f"notes: {text_frame.text}")
                    slide.notes_slide.notes_text_frame.text = text_frame.text
                    # response = translate.translate_text(
                    #         Text=text_frame.text,
                    #         SourceLanguageCode=source_language_code,
                    #         TargetLanguageCode=target_language_code,
                    #         TerminologyNames=terminology_names)
                    # slide.notes_slide.notes_text_frame.text = response.get('TranslatedText')
                    # response_str = translate_from_ollama()
                except Exception as err:
                    print(f"Unexpected {err=}, {type(err)=}")

        # translate other texts
        for shape in slide.shapes:
            translate_shape(shape, source_language_code, target_language_code, terminology_names)


# def import_terminology(terminology_file_path):
#     print('Importing terminology data from {file_path}...'.format(file_path=terminology_file_path))
#     with open(terminology_file_path, 'rb') as f:
#         translate.import_terminology(Name=TERMINOLOGY_NAME,
#                                      MergeStrategy='OVERWRITE',
#                                      TerminologyData={'File': bytearray(f.read()), 'Format': 'CSV'})


def main():
    argument_parser = argparse.ArgumentParser(
            'Translates pptx files from source language to target language using LLM')
    argument_parser.add_argument(
            'source_language_code', type=str,
            help='The language code for the language of the source text. Example: en')
    argument_parser.add_argument(
            'target_language_code', type=str,
            help='The language code requested for the language of the target text. Example: pt')
    argument_parser.add_argument(
            'input_file_path', type=str,
            help='The path of the pptx file that should be translated')
    # argument_parser.add_argument(
    #         '--terminology', type=str,
            # help='The path of the terminology CSV file')
    args = argument_parser.parse_args()

    terminology_names = []
    # if args.terminology:
    #     import_terminology(args.terminology)
    #     terminology_names = [TERMINOLOGY_NAME]

    # translate_from_ollama("hello world!", 
    #                         args.source_language_code,
    #                         args.target_language_code)
    
    get_dont_translate_words(DONT_TRANSLATE_WORDS_FILE)
    print('Translating {file_path} from {source_language_code} to {target_language_code}...'.format(
            file_path=args.input_file_path,
            source_language_code=args.source_language_code,
            target_language_code=args.target_language_code))
    presentation = Presentation(args.input_file_path)
    translate_presentation(presentation,
                           args.source_language_code,
                           args.target_language_code,
                           terminology_names)

    output_file_path = args.input_file_path.replace(
            '.pptx', '-{language_code}.pptx'.format(language_code=args.target_language_code))
    print('Saving {output_file_path}...'.format(output_file_path=output_file_path))
    presentation.save(output_file_path)


if __name__== '__main__':
    main()
