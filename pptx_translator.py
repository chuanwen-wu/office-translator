#!/usr/bin/env python
import logging
import argparse
from pptx import Presentation
from pptx.text.text import TextFrame
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.shapes.autoshape import Shape
import json
import base64
import requests
from pptx.util import Pt
from io import BytesIO
import os
import traceback

KAFKA_CONFIG = {
    'topic_pptx': os.getenv('KAFKA_TOPIC_PPTX', 'pptx-translate'),
    'topic_status': os.getenv('KAFKA_TOPIC_STATUS_UPDATE', 'status-update'),
    'servers': os.getenv('KAFKA_SERVERS', 'localhost:9092').split(',')
}

OLLAMA_CONFIG = {
    'url': os.getenv('OLLAMA_URL', 'http://localhost:11434')
}

"""
Dict that maps language code to MSO_LANGUAGE_ID enum value.
"""
LANGUAGE_CODE_TO_LANGUAGE_ID = {
    'en': MSO_LANGUAGE_ID.ENGLISH_US,
    'zh': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    # 'zh-HK': MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
    # 'zh-TW': MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
}

# 这里跳过，原因是subscribe_translate_task 如果连接kafka异常，以抛异常退出
def check_kafka():
    return True

def check_ollama():
    logger.info(f"OLLAMA_URL: {OLLAMA_CONFIG['url']}")
    try:
        resp = requests.head(OLLAMA_CONFIG['url'])
        logger.debug(f"status_code: {resp.status_code}")
        logger.debug(f"text: {resp.text}")
        logger.debug(f"content: {resp.content}")
        logger.debug(f"header: {resp.headers}")
        if resp.status_code == 200:
            logger.info("running normally: head return 200")
            return True
        else:
            raise Exception(f'ollama resp error, status_code={resp.status_code}')
    except Exception as err:
        logger.error(f"Unexpected {err=}, {type(err)=}")
        logger.error(f"connect ollama endpoint failed: {OLLAMA_CONFIG['url']}")
        return False

DONT_TRANSLATE_WORDS_FILE = f"./dont-translate-word.txt"
DONT_TRANSLATE_WORDS = f""

# Set up logger
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)
formatter = logging.Formatter("%(asctime)s - %(filename)s:%(lineno)d - %(levelname)s - %(message)s")

# Log to console
handler = logging.StreamHandler()
handler.setFormatter(formatter)
logger.addHandler(handler)

def contain_chinese(string):
    """
    检查整个字符串是否包含中文
    :param string: 需要检查的字符串
    :return: bool
    """
    for ch in string:
        if u'\u4e00' <= ch <= u'\u9fff':
            return True
 
    return False

def get_dont_translate_words(file_path):
    global DONT_TRANSLATE_WORDS
    arr = []
    with open(file_path, "r") as f:
        text = f.read()
    
    return parse_dont_translate_words(text)

'''
输入text是一段以换行符作为分割符的文本，比如：
    AWS
    CDN
输出:
    AWS,CDN
'''
def parse_dont_translate_words(text: str):
    global DONT_TRANSLATE_WORDS
    arr = []
    for l in text.splitlines():
        if l.strip().startswith('#'):
            continue
        arr.append(l.strip())
    
    DONT_TRANSLATE_WORDS = ", ".join(arr)
    return DONT_TRANSLATE_WORDS

def translate_from_ollama(src, source_language_code, target_language_code, model:str="qwen2.5", temperature:int=0):
    source_language = ""
    target_language = ""
    if source_language_code.lower() == 'en':
        source_language = "English"
    elif source_language_code.lower() == 'zh':
        source_language = "Chinese"

    if target_language_code.lower() == 'en':
        target_language = "English"
    elif target_language_code.lower() == 'zh':
        target_language = "Chinese"

    temp = {
        "model": model, # "llama3.2" or "qwen2.5"
        "stream": False,
        "options": {
            'temperature': temperature
        },
        "messages": [
            {
                "role": "system",
                "content": f"You are a highly skilled translator. Your task is to accurately translate the text I provide from \
{source_language} into the {target_language} while preserving the meaning, tone, and nuance of the original text. \
Please maintain proper grammar, spelling, and punctuation in the translated version. \
Please answer briefly without any other hints or extensions. \
If the input text does not make sense or empty or incomplete, just return the original text. \
Remove the leading and trailing double quotes. \
Please keeping these phrases unchanged: {DONT_TRANSLATE_WORDS}."
            },
            {"role": "user", "content": f"{src}"}
        ]
    }
    data = json.dumps(temp)
    logger.info(f'temp: {data}')
    # url = 'http://localhost:11434/api/generate'
    # url = 'http://localhost:11434/api/chat'
    url = f"{OLLAMA_CONFIG['url']}/api/chat"
    resp = requests.post(url, data=data)
    if resp.status_code == 200:
        body = json.loads(resp.text)
        logger.info(f"{src} --> {body['message']['content']}")
        return body['message']['content']
    else: # http code != 200
        raise Exception(f'ollama resp error, status_code={resp.status_code}, text={resp.text}')

'''
根据文本框的大小，自动调整文本的字体大小
'''
def auto_fit_text(text_frame: TextFrame, max_font_size=40, min_font_size=6, font_family='Arial', text_original='', font_size_original:float=12.0) -> int:     
    logger.debug(f"[auto_fit_text]: max_font_size={max_font_size}, min_font_size={min_font_size}, font_size_original={font_size_original}")
    def emu_to_cm(emu):
        return round(emu / 360000, 2)
    
    # option 2: try to resize the font size by the text frame size
    fs = get_font_size_by_text_rect(text_frame, font_family, max_font_size, False, False)
    logger.debug(f"get_font_size_by_text_rect = {fs}")
    if fs > min_font_size and fs < max_font_size:
        logger.debug(f"option2: set font size = {fs}")
        text_frame.word_wrap = True
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(fs)
    else:   
        try:
            logger.debug(text_frame._parent) # pptx.shapes.autoshape.Shape object
            logger.debug(f"Width: {emu_to_cm(text_frame._parent.width)} cm, Height: {emu_to_cm(text_frame._parent.height)} cm")
        except Exception as err:
            logger.error(f"Unexpected {err=}, {type(err)=}")
            logger.error(traceback.format_exc())
        # option 3: try to resize the font size by the original text size
        if font_size_original > 0:
            fs = get_font_size_by_text_original(text_original, text_frame.text, font_size_original)
            if fs < min_font_size:
                fs = min_font_size
            if fs > max_font_size:
                fs = max_font_size
            logger.debug(f"option3: set font size = {fs}")
            text_frame.word_wrap = True
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(fs)
        else:
            logger.debug("font_size_original is 0, do nothing")
        

# option 2
def get_font_size_by_text_rect(text_frame: TextFrame, font_family='Arial', max_font_size=18, bold=False, italic=False, font_file='assets/fonts/Arial Unicode.ttf'):
    if type(text_frame._parent) is pptx.table._Cell:
        logger.debug("text_frame._parent is table cell, can't auto fit text")
        return 0
    
    for i in range(max_font_size, 0, -1):
        # logger.info(f"try {i} ------------")
        try:
            # text_frame.fit_text(font_family=font_family, max_size=i, bold=bold, italic=italic, font_file=font_file)
            # text_frame.word_wrap = True
            # logger.info(f"fit font size = {i}")
            # return 0
            fs = text_frame._best_fit_font_size(font_family, i, bold, italic, font_file)
            if fs is None:
                fs = 0
            return fs
        except TypeError as err:
            # logger.debug(f"{err}")
            # logger.debug(traceback.format_exc())
            pass  # ignore TypeError
        except Exception as err:
            logger.error(f"Unexpected {err=}, {type(err)=}")
            logger.error(traceback.format_exc())
            break  # break the loop
    logger.error('Error: could not find fit font size!')
    return 0

from PIL import Image, ImageDraw, ImageFont
import math
def get_font_size_by_text_original(text_orininal: str, text_translated: str, font_size_original: float) -> float:
    def _rendered_size(text, point_size, font_file):
        """
        Return a (width, height) pair representing the size of *text* in English
        Metric Units (EMU) when rendered at *point_size* in the font defined in
        *font_file*.
        """
        emu_per_inch = 914400
        px_per_inch = 72.0
        font = ImageFont.truetype(font_file, point_size)
        try:
            px_width, px_height = font.getsize(text)
        except AttributeError:
            left, top, right, bottom = font.getbbox(text)
            px_width, px_height = right - left, bottom - top

        emu_width = int(px_width / px_per_inch * emu_per_inch)
        emu_height = int(px_height / px_per_inch * emu_per_inch)
        return emu_width, emu_height
    
    font_path = "./assets/fonts/Arial Unicode.ttf"
    len_original = len(text_orininal)
    len_translated = len(text_translated)

    w1, h1 = _rendered_size(text_orininal, font_size_original, font_path)
    logger.debug(f"{text_orininal}, {font_size_original}, {w1}, {h1}")

    w2, h2 = _rendered_size(text_translated, font_size_original, font_path)
    logger.debug(f"{text_translated}, {font_size_original}, {w2}, {h2}")
    
    # w1*h1 = w2*h2*n^2
    # option 1:
    # fs = round(math.sqrt(w1/w2) * font_size_original, 1)

    # option 2:
    fs = round(math.sqrt(w1*h1/(h2*w2)) * font_size_original, 1)

    w3, h3 = _rendered_size(text_translated, fs, font_path)
    logger.debug(f"{text_translated}, {fs}, {w3}, {h3}")

    return fs

def translate_text_frame(text_frame: TextFrame, source_language_code, target_language_code, terminology_names, auto_resize_text=True):
    def delete_run(run):
        r = run._r
        r.getparent().remove(r)

    changed = False
    original_max_fs_pt = None
    original_min_fs_pt = None
    text_original = text_frame.text
    for paragraph in text_frame.paragraphs:
        paraText = ''
        for index, paragraph_run in enumerate(paragraph.runs):
            paraText += paragraph_run.text
            
        try:
            paraText = paraText.strip()
            if len(paraText) == 0:
                continue

            # 源语言是中文，但实际文本未包含中文，跳过
            if contain_chinese(paraText.strip()) == False and source_language_code.lower() == 'zh':
                logger.info(f"not chinese: {paraText}")
                continue

            response_text = translate_from_ollama(paraText, source_language_code, target_language_code)
            paragraph.runs[0].text = response_text
            paragraph.runs[0].font.language_id = LANGUAGE_CODE_TO_LANGUAGE_ID[target_language_code]
            logger.info(f"font: id={paragraph.runs[0].font.language_id}, name={paragraph.runs[0].font.name}, size(pt)={paragraph.runs[0].font.size.pt if paragraph.runs[0].font.size is not None else None}")
            size = len(paragraph.runs)
            index = size - 1
            while index > 0:
                delete_run(paragraph.runs[index])
                index = index - 1

            # 取得原本段落中最大的字体，以及最小的字体
            if paragraph.runs[0].font.size is not None:
                if original_max_fs_pt is None:
                    original_max_fs_pt = paragraph.runs[0].font.size.pt
                elif paragraph.runs[0].font.size.pt > original_max_fs_pt:
                    original_max_fs_pt = paragraph.runs[0].font.size.pt
                if original_min_fs_pt is None:
                    original_min_fs_pt = paragraph.runs[0].font.size.pt
                elif paragraph.runs[0].font.size.pt < original_min_fs_pt:
                    original_min_fs_pt = paragraph.runs[0].font.size.pt

            changed = True
        except Exception as err:
            logger.error(f"Unexpected {err=}, {type(err)=}")
            raise err

    if auto_resize_text and changed and len(text_frame.text.strip()) > 0:
        font_size_original = (original_max_fs_pt+original_min_fs_pt)/2 if original_max_fs_pt is not None and original_min_fs_pt is not None else 0
        if original_max_fs_pt is None:
            original_max_fs_pt = 40
        if original_min_fs_pt is None:
            original_min_fs_pt = 6
        if source_language_code.lower() == 'zh' and target_language_code.lower() == 'en':
            auto_fit_text(text_frame, max_font_size=int(original_max_fs_pt), min_font_size=6, 
                          text_original=text_original, font_size_original=font_size_original)
            # text_frame.word_wrap = True
            # text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        elif source_language_code.lower() == 'en' and target_language_code.lower() == 'zh':
            auto_fit_text(text_frame, max_font_size=40, min_font_size=int(original_min_fs_pt), 
                          text_original=text_original, font_size_original=font_size_original)

def translate_table(table, source_language_code, target_language_code, terminology_names):
    for row in table.rows:
        for cell in row.cells:
            text_frame = cell.text_frame
            translate_text_frame(text_frame, source_language_code, target_language_code, terminology_names)
            

def translate_shape(shape, source_language_code, target_language_code, terminology_names, auto_resize_text=True):
    # table
    if shape.has_table:
        translate_table(shape.table, source_language_code, target_language_code, terminology_names)

    # text_frame
    if shape.has_text_frame:
        translate_text_frame(shape.text_frame, source_language_code, target_language_code, terminology_names, auto_resize_text=auto_resize_text)

    # groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            translate_shape(sub_shape, source_language_code, target_language_code, terminology_names, auto_resize_text=auto_resize_text)
    

def translate_presentation(presentation: Presentation, source_language_code, target_language_code, terminology_names, auto_resize_text=True):
    slide_number = 1
    for slide in presentation.slides:
        logger.info('Slide {slide_number} of {number_of_slides}'.format(
                slide_number=slide_number,
                number_of_slides=len(presentation.slides)))
        slide_number += 1

        # 尚不支持smartArt
        try:
            # translate comments
            if slide.has_notes_slide:
                text_frame = slide.notes_slide.notes_text_frame
                if len(text_frame.text) > 0:
                    logger.info(f"notes: {text_frame.text}")
                    # slide.notes_slide.notes_text_frame.text = text_frame.text
                    # response = translate.translate_text(
                    #         Text=text_frame.text,
                    #         SourceLanguageCode=source_language_code,
                    #         TargetLanguageCode=target_language_code,
                    #         TerminologyNames=terminology_names)
                    # slide.notes_slide.notes_text_frame.text = response.get('TranslatedText')
                    response_str = translate_from_ollama(slide.notes_slide.notes_text_frame.text, source_language_code, target_language_code)
                    slide.notes_slide.notes_text_frame.text = response_str

            # translate body
            for shape in slide.shapes:
                translate_shape(shape, source_language_code, target_language_code, terminology_names, auto_resize_text=auto_resize_text)
        except Exception as err:
            logger.error(f"Unexpected err: {err}, {type(err)}")
            raise err
    
def run_as_local():
    logger.info("Run in standalone mode")
    argument_parser = argparse.ArgumentParser(
        'Translates pptx files using LLM in local mode')
    argument_parser.add_argument(
        '-s', '--standalone',
        action="store_true",
        help='Run as local script')
    argument_parser.add_argument(
        '--source_language_code', '--sl', type=str,
        required=True,
        help='The language code for the language of the source text. Example: en')
    argument_parser.add_argument(
        '--target_language_code', '--tl', type=str,
        required=True,
        help='The language code requested for the language of the target text. Example: zh')
    argument_parser.add_argument(
        '--input_file_path', '--if', type=str,
        required=True,
        help='The path of the pptx file that should be translated')
    args = argument_parser.parse_args()

    terminology_names = []
    
    get_dont_translate_words(DONT_TRANSLATE_WORDS_FILE)
    logger.info('Translating {file_path} from {source_language_code} to {target_language_code}...'.format(
            file_path=args.input_file_path,
            source_language_code=args.source_language_code,
            target_language_code=args.target_language_code))
    if not check_language_code(args.source_language_code):
        logger.error(f"Invalid source_language_code: {args.source_language_code}")
        return 1, "Invalid source_language_code"
    if not check_language_code(args.target_language_code):
        logger.error(f"Invalid target_language_code: {args.target_language_code}")
        return 1, "Invalid target_language_code"
    
    presentation = Presentation(args.input_file_path)
    translate_presentation(presentation,
                           args.source_language_code,
                           args.target_language_code,
                           terminology_names)

    output_file_path = args.input_file_path.replace(
            '.pptx', '-{language_code}.pptx'.format(language_code=args.target_language_code))
    logger.info('Saving {output_file_path}...'.format(output_file_path=output_file_path))
    presentation.save(output_file_path)


from kafka import KafkaConsumer
from time import sleep
def subscribe_translate_task():
    topic_name = KAFKA_CONFIG['topic_pptx'] #'pptx-translate'
    group_name = 'group1'
    kafka_servers = KAFKA_CONFIG['servers'] #['localhost:9092']
    logger.info(f"[subscribe_translate_task] Kafka config: {KAFKA_CONFIG}")
    while True:
        # To consume latest messages and auto-commit offsets
        consumer = KafkaConsumer(topic_name,
                                 group_id=group_name,
                                 bootstrap_servers=kafka_servers,
                                 enable_auto_commit=True,
                                #  auto_offset_reset='earliest',
                                 max_poll_records=1)
        logger.info(f"Waiting for task...")
        for message in consumer:
            # message value and key are raw bytes -- decode if necessary!
            # e.g., for unicode: `message.value.decode('utf-8')`
            # print("%s:%d:%d: key=%s value=%s" % (message.topic, message.partition,
            #                                   message.offset, message.key,
            #                                   message.value))
            #logger.info("%s:%d:%d: key=%s" % (message.topic, message.partition, message.offset, message.key))
            task = json.loads(message.value.decode('utf-8'))
            # task['input_file_content'] = base64.b64decode(task['input_file_content'])
            logger.info(f"Receive task: id={task['id']}, md5={task['md5']}, source_language={task['source_language']}, target_language={task['target_language']}, input_file_path={task['input_file_path']}, output_file_path={task['output_file_path']}, auto_resize_text={task['auto_resize_text']}")
            break
        consumer.close()
        # task['status'] = 1 #processing
        publish_status_update({
            'id': task['id'],
            'status': 1,
            'source_language': task['source_language'],
            'target_language': task['target_language'],
            'md5': task['md5'],
            'auto_resize_text': task['auto_resize_text']
        })

        if 'dont_translate_list' in task and task['dont_translate_list'] != '':
            parse_dont_translate_words(task['dont_translate_list'])
        else:
            get_dont_translate_words(DONT_TRANSLATE_WORDS_FILE)

        try: 
            with open(task['input_file_path'], 'rb') as f:
                input_file_content = f.read()
            ret, output_file_content = translate_pptx_file(input_file_content, task['source_language'], task['target_language'], auto_resize_text=task['auto_resize_text'])
            if ret == 0:
                logger.info(f"write translated file to {task['output_file_path']}")
                with open(task['output_file_path'], 'wb') as f:
                    f.write(output_file_content)
                task['status'] = 2  #finished                
            else: #error
                err = output_file_content
                logger.error(f"translate_pptx_file failed: ret={ret}, err={err}")
                task['status'] = 3
                task['error_msg'] = err
        except Exception as err2:
            logger.error(err2)
            task['status'] = 3
            task['error_msg'] = str(err2)
        # del task['input_file_content']
        publish_status_update(task)

from kafka import KafkaProducer
from kafka.errors import KafkaError
def publish_status_update(task):
    logger.info(f"[publish_status_update]: id={task['id']}, status={task['status']}")

    topic_name = KAFKA_CONFIG['topic_status'] #'status-update'
    kafka_servers = KAFKA_CONFIG['servers'] #['localhost:9092']
    producer = KafkaProducer(bootstrap_servers=kafka_servers, max_request_size=1024*1024*1024)
    # if task['status'] == 2:
    #     task['output_file_content'] = base64.b64encode(task['output_file_content']).decode('utf-8')
    # Asynchronous by default
    future = producer.send(topic_name, key=bytes(str(task['id']), encoding='ascii'), value=json.dumps(task).encode('utf-8'))
    # Block for 'synchronous' sends
    try:
        record_metadata = future.get(timeout=10)
        # Successful result returns assigned partition and offset
        logger.info(f"record_metadata.topic: {record_metadata.topic}, partition: {record_metadata.partition}, offset: {record_metadata.offset}")
        return True
    except KafkaError as err:
        logger.info(f'KafkaError: {err}')
        pass

def check_language_code(language_code):
    return True if language_code in LANGUAGE_CODE_TO_LANGUAGE_ID else False

# return (ret, output_file_content)
# ret =0 表示成功；
# ret=1表示输入文件有错误；
# ret=2表示翻译过程出错
def translate_pptx_file(input_file_content: bytes, source_language_code, target_language_code, auto_resize_text=True):
    ret = 0
    logger.info(f"[translate_pptx_file] source_language_code: {source_language_code}, target_language_code: {target_language_code}")
    
    if not check_language_code(source_language_code):
        logger.error(f"Invalid source_language_code: {source_language_code}")
        return 1, "Invalid source_language_code"
    if not check_language_code(target_language_code):
        logger.error(f"Invalid target_language_code: {target_language_code}")
        return 1, "Invalid target_language_code"
    try:
        presentation = Presentation(BytesIO(input_file_content))
    except Exception as err:
        ret = 1
        logger.error(f"Unexpected {err}, {type(err)}")
        return ret, str(err)
    try:
        translate_presentation(presentation, source_language_code, target_language_code, None, auto_resize_text=auto_resize_text)
        output_file_content = BytesIO()
        presentation.save(output_file_content)
        return ret, output_file_content.getvalue()
    except Exception as err:
        ret = 2
        logger.error(f"Unexpected {err}, {type(err)}")
        return ret, str(err)


def init_check():
    if not check_ollama():
        exit(1)
    if not check_kafka():
        exit(2)

def test():
    try:
        with open("scripts/docker/data/file_repo/71d6ede295b2d7f0ec3c3f9df7d16a27-input.pptx", 'rb') as f:
            input_file_content = f.read()
    except Exception as err:
        logger.error(err)
        return 1

    ret, output_file_content = translate_pptx_file(input_file_content, 'en', 'zh')
    logger.info(f"ret={ret}")

    try:
        with open("scripts/docker/data/file_repo2/output.pptx", 'wb') as f:
            f.write(output_file_content)
    except Exception as err:
        logger.error(err)
        return 2

    logger.info(f"done")

def run_as_service():
    logger.info("Run in service mode")
    init_check()
    subscribe_translate_task()

def init():
    argument_parser = argparse.ArgumentParser(
            'Translates pptx files from source language to target language using LLM')
    argument_parser.add_argument(
            '-s', '--standalone',
            action="store_true",
            help='Run as local script. Default is service mode')
    argument_parser.add_argument(
            '--sl', '--source_language_code', type=str,
            help='The language code for the language of the source text. Example: en')
    argument_parser.add_argument(
            '--tl', '--target_language_code', type=str,
            help='The language code requested for the language of the target text. Example: zh')
    argument_parser.add_argument(
            '--if', '--input_file_path', type=str,
            help='The path of the pptx file that should be translated')
    args = argument_parser.parse_args()
    if args.standalone:
        run_as_local()
    else:
        run_as_service()

import pptx
if __name__== '__main__':
    logger.info(f"OLLAMA_CONFIG: {OLLAMA_CONFIG}")
    # print(pptx.version)
    init()
